import sys
import pandas as pd
import json
import ast
from pathlib import Path
from datetime import datetime
import streamlit as st
import warnings
import altair as alt
import base64
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

# --- Path Correction ---
project_root = Path(__file__).resolve().parents[1]
sys.path.append(str(project_root))
# -----------------------

warnings.filterwarnings("ignore", message=".*Pydantic BaseModel V1.*")

from backend.graph.build_graph import build_graph
from backend.agents.chatbot_agent import run_chatbot_agent
from backend.agents.drift_linker_agent import run_drift_linker_agent
from backend.utils.ingest_documents import process_context_files
from backend.agents.drift_linker_agent import ConnectionType

# --- Descriptions for Drift Linker Categories ---
CONNECTION_TYPE_DESCRIPTIONS = {
    ConnectionType.STRONG_CAUSAL.value: "One drift appears to be a direct cause of another subsequent drift.",
    ConnectionType.SHARED_EVIDENCE.value: "The drifts are linked by common evidence or appear to share the same underlying root cause.",
    ConnectionType.THEMATIC_OVERLAP.value: "The drifts are not directly linked but share a similar theme or occur in related business areas.",
}

# --- UI Helper Function ---
def get_date_from_filename(filename: str) -> str:
    """Parses a YYYY-MM-DD date from the start of a filename."""
    try:
        filename = Path(filename).name
        date_str = filename.split('_')[0]
        dt_obj = datetime.strptime(date_str, '%Y-%m-%d')
        return dt_obj.strftime('%d.%m.%Y')
    except (ValueError, IndexError):
        return "N/A"

# --- Helper Function to Load and Unpack Drifts ---
@st.cache_data
def load_and_unpack_drifts(data_dir_path: str):
    """
    Loads drift data from a specific event log directory.
    Returns a tuple: (drift_options, error_message).
    """
    drift_dir = Path(data_dir_path)
    if not drift_dir.exists():
        return [], f"Could not find drift data directory at {drift_dir}"

    try:
        csv_path = next(drift_dir.glob("*.csv"))
        df = pd.read_csv(csv_path)
        
        drift_options = []
        for row_index, row in df.iterrows():
            drift_types = ast.literal_eval(row['Detected Drift Types'])
            changepoints = ast.literal_eval(row['Detected Changepoints'])
            gold_docs = ast.literal_eval(row.get('gold_source_document', '[]'))
            
            for drift_index, drift_type in enumerate(drift_types):
                option_id = f"{row_index}-{drift_index}"
                display_name = (
                    f"Drift #{len(drift_options) + 1}: {drift_type} "
                    f"(between {changepoints[drift_index][0]} and {changepoints[drift_index][1]})"
                )
                drift_options.append({
                    "id": option_id, 
                    "display": display_name, 
                    "data_dir": str(drift_dir),
                    "gold_doc": gold_docs[drift_index] if drift_index < len(gold_docs) else None
                })
        
        return drift_options, None
        
    except StopIteration:
        return [], f"No CSV file found in the drift data directory: {drift_dir}"
    except Exception as e:
        return [], f"Error loading or parsing drift data: {e}"

# --- Visualization Helper ---
def create_timeline_chart(drift_info: dict, ranked_causes: list):
    """Creates an Altair timeline chart for a single drift."""
    try:
        start_date = pd.to_datetime(drift_info.get("start_timestamp"))
        end_date = pd.to_datetime(drift_info.get("end_timestamp"))

        drift_period_df = pd.DataFrame([{
            "start": start_date,
            "end": end_date,
            "event": "Drift Period"
        }])

        evidence_points = []
        for cause in ranked_causes:
            # The 'timestamp' is a Unix timestamp
            doc_timestamp = cause.get("timestamp", 0)
            if doc_timestamp > 0:
                doc_date = pd.to_datetime(doc_timestamp, unit='s')
                doc_name = Path(cause.get("source_document")).name
                evidence_points.append({"date": doc_date, "event": doc_name})
        
        evidence_df = pd.DataFrame(evidence_points)

        drift_bar = alt.Chart(drift_period_df).mark_bar(
            opacity=0.3,
            color='#FF4B4B'
        ).encode(
            x=alt.X('start:T', title="Timeline"),
            x2='end:T',
            tooltip=[alt.Tooltip('start:T', title='Drift Start'), alt.Tooltip('end:T', title='Drift End')]
        )

        evidence_marks = alt.Chart(evidence_df).mark_point(
            size=100,
            filled=True,
            color='#33C481',
            opacity=0.9
        ).encode(
            x='date:T',
            # --- CHANGE START ---
            # Explicitly define the data type for the 'event' field as Nominal (string)
            tooltip=[alt.Tooltip('event:N', title='Document'), alt.Tooltip('date:T', title='Published Date')]
            # --- CHANGE END ---
        )
        
        return (drift_bar + evidence_marks).properties(height=100)
    except Exception as e:
        st.warning(f"Could not generate timeline chart: {e}")
        return None
    
# --- Document Viewer Helper ---
def display_document(file_path_str: str):
    file_path = Path(file_path_str)
    if not file_path.exists():
        st.error(f"File not found: {file_path}")
        return

    st.subheader(f"Preview: {file_path.name}")
    
    if file_path.suffix.lower() == ".pdf":
        with open(file_path, "rb") as f:
            base64_pdf = base64.b64encode(f.read()).decode('utf-8')
        pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800" type="application/pdf"></iframe>'
        st.markdown(pdf_display, unsafe_allow_html=True)
    elif file_path.suffix.lower() == ".pptx":
        with st.spinner("Converting PowerPoint slides..."):
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                st.write(f"--- Slide {i + 1} ---")
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                st.text("\n".join(slide_text))
    else:
        st.info(f"Preview for {file_path.suffix} files is not yet supported. Showing plain text content.")
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                st.text(f.read())
        except Exception as e:
            st.error(f"Could not read file as text: {e}")

# --- Page Configuration ---
st.set_page_config(page_title="Concept Drift Explainer", page_icon="🤖", layout="wide")

# --- Session State Initialization ---
# Initialize each key individually to ensure they exist.
if 'analysis_run_complete' not in st.session_state:
    st.session_state.analysis_run_complete = False
if 'all_explanations' not in st.session_state:
    st.session_state.all_explanations = []
if 'error_message' not in st.session_state:
    st.session_state.error_message = None
if 'feedback_states' not in st.session_state:
    st.session_state.feedback_states = {}
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'full_state_log' not in st.session_state:
    st.session_state.full_state_log = []
if 'linked_drift_summary' not in st.session_state:
    st.session_state.linked_drift_summary = None
if 'connection_type' not in st.session_state:
    st.session_state.connection_type = None
if 'show_chat' not in st.session_state:
    st.session_state.show_chat = False
if 'selected_log' not in st.session_state:
    st.session_state.selected_log = None
if 'doc_to_view' not in st.session_state:
    st.session_state.doc_to_view = None
if 'dark_mode' not in st.session_state:
    st.session_state.dark_mode = True
if 'active_tab' not in st.session_state:
    st.session_state.active_tab = "Main Analysis"

# === THEME INJECTION ===
if st.session_state.dark_mode:
    st.markdown(
        """
        <style>
          body {background-color:#0E1117; color:#E4E6EB;}
          .sidebar .css-1d391kg {background-color:#121620;}
        </style>
        """, unsafe_allow_html=True
    )
else:
    st.markdown(
        """
        <style>
          body {background-color:#FFFFFF; color:#000000;}
          .sidebar .css-1d391kg {background-color:#F5F5F5;}
        </style>
        """, unsafe_allow_html=True
    )
# ========================

st.title("Concept Drift Explainer")

# --- Sidebar ---
with st.sidebar:
    st.header("Controls")
    
    event_logs_dir = project_root / "data" / "event_logs"
    event_log_options = {d.name: str(d) for d in event_logs_dir.iterdir() if d.is_dir()}
    
    if st.session_state.selected_log is None:
            st.session_state.selected_log = list(event_log_options.values())[0]

    selected_log_name = st.selectbox(
        "Select Event Log",
        options=list(event_log_options.keys()),
        index=list(event_log_options.values()).index(st.session_state.selected_log)
    )

    if st.session_state.selected_log != event_log_options[selected_log_name]:
        st.session_state.selected_log = event_log_options[selected_log_name]
        st.cache_data.clear()
        st.rerun()

    drift_options, load_error = load_and_unpack_drifts(st.session_state.selected_log)

    if drift_options:
        st.info(f"Found **{len(drift_options)}** drift(s) to analyze.")
        if st.button("Run Full Analysis"):
            with st.spinner(f"Analyzing {len(drift_options)} drift(s)..."):
                try:
                    app = build_graph()
                    log_explanations = []
                    log_states = []
                    for drift in drift_options:
                        row_idx, drift_idx = map(int, drift['id'].split('-'))
                        initial_input = {"selected_drift": {"row_index": row_idx, "drift_index": drift_idx, "data_dir": drift["data_dir"]}}
                        final_state = app.invoke(initial_input)
                        if final_state.get("error"):
                            st.session_state.error_message = f"Error: {final_state['error']}"
                            break
                        log_explanations.append(final_state.get('explanation'))
                        log_states.append(final_state)
                    if not st.session_state.error_message:
                        st.session_state.all_explanations.extend(log_explanations)
                        st.session_state.full_state_log.extend(log_states)
                        st.session_state.analysis_run_complete = True
                        if len(log_states) > 1:
                            linker_result = run_drift_linker_agent(log_states)
                            st.session_state.linked_drift_summary = linker_result.get("linked_drift_summary")
                            st.session_state.connection_type = linker_result.get("connection_type")
                except Exception as e:
                    st.session_state.error_message = f"Analysis failed: {e}"
            st.rerun()
    elif load_error:
        st.error(load_error)
    else:
        st.info("No drifts were detected in this log.")

    # Document Upload field
    DOCUMENTS_PATH = project_root / "data" / "documents"
    st.markdown("---")
    st.header("Add New Context")
    uploaded_files = st.file_uploader(
        "Upload new documents (.pdf, .pptx, etc.)",
        accept_multiple_files=True,
        type=['pdf', 'pptx', 'docx', 'txt', 'png', 'jpg']
    )

    if st.button("Process Uploaded Files") and uploaded_files:
        saved_files_paths = []
        with st.spinner("Saving and processing uploaded files..."):
            for uploaded_file in uploaded_files:
                # We must save the file to our documents folder first
                file_path = DOCUMENTS_PATH / uploaded_file.name
                with open(file_path, "wb") as f:
                    f.write(uploaded_file.getbuffer())
                saved_files_paths.append(file_path)
            
            # Now, call the ingestion logic with only the new files
            process_context_files(saved_files_paths)
            st.success(f"Successfully processed {len(saved_files_paths)} new document(s)!")
            # Clear the drift data cache so the UI reloads it if necessary
            st.cache_data.clear()
            st.rerun()

    st.markdown("---")
    st.header("Settings")
    st.checkbox("Dark Mode", value=True, key="dark_mode")

# --- Callback for Tab Switching ---
def set_active_tab(tab_name, doc_path=None):
    st.session_state.active_tab = tab_name
    if doc_path:
        st.session_state.doc_to_view = doc_path

# --- Main Content Area ---
tab_options = ["Main Analysis", "Document Viewer"]
active_tab_index = tab_options.index(st.session_state.active_tab)
selected_tab = st.radio("Navigation", tab_options, index=active_tab_index, key="tab_selector", horizontal=True)

if selected_tab != st.session_state.active_tab:
    set_active_tab(selected_tab)
    st.rerun()

if st.session_state.active_tab == "Main Analysis":
    if st.session_state.error_message:
        st.error(st.session_state.error_message)
    elif st.session_state.analysis_run_complete:
        if not st.session_state.all_explanations:
             st.info("Analysis complete. No drifts were detected in this log.")
        else:
            st.success(f"Successfully analyzed {len(st.session_state.all_explanations)} drift(s).")
        
        if st.session_state.linked_drift_summary:
            with st.container():
                st.subheader("🔗 Cross-Drift Analysis")
                desc = CONNECTION_TYPE_DESCRIPTIONS.get(st.session_state.connection_type, "")
                st.markdown(f"**Connection:** {desc}")
                st.markdown(st.session_state.linked_drift_summary)
        
        for idx, explanation in enumerate(st.session_state.all_explanations):
            drift_state = st.session_state.full_state_log[idx]
            drift_info = drift_state.get("drift_info", {})
            st.subheader(f"Explanation for Drift #{idx + 1}: {drift_info.get('drift_type', 'Unknown').capitalize()}")
            
            ranked_causes = explanation.get('ranked_causes', [])

            # Generate and display the timeline chart
            timeline_chart = create_timeline_chart(drift_info, ranked_causes)
            if timeline_chart:
                st.altair_chart(timeline_chart, use_container_width=True)

            for cause_idx, cause in enumerate(explanation.get('ranked_causes', [])):
                with st.container():
                    cols = st.columns([1, 2, 2])
                    with cols[0]:
                        st.metric("Drift Trigger Date", get_date_from_filename(cause.get("source_document", "")))
                    with cols[1]:
                        st.markdown(f"**Source Document:**")
                        if st.button(cause.get('source_document', 'N/A'), key=f"doc_link_{idx}_{cause_idx}"):
                            st.session_state.doc_to_view = cause.get('source_document')
                            # A simple st.rerun() will not switch tabs. A more complex callback is needed.
                            # For now, we instruct the user.
                            st.info("Document viewer tab is ready.")
                    with cols[2]:
                        st.markdown(f"**Confidence:** `{cause.get('confidence_score', 0.0)*100:.1f}%`")

                    with st.expander("View Details"):
                        # The full details and feedback logic are now included
                        st.markdown(f"**Description:** {cause.get('cause_description', 'N/A')}")
                        st.markdown("**Evidence:**")
                        st.code(cause.get('evidence_snippet', 'N/A'), language='text')
                        st.caption(f"Source Document: `{cause.get('source_document', 'N/A')}`")

                    # Granular Feedback Mechanism
                    st.markdown("---")
                    feedback_key = (idx, cause_idx)
                    if st.session_state.feedback_states.get(feedback_key):
                            st.success("Thank you for your feedback on this cause!")
                    else:
                        st.write("Was this specific cause helpful?")
                        col1, col2, _ = st.columns([1, 1, 8])
                        with col1:
                            # Use both explanation_idx and cause_idx for a unique key
                            if st.button("👍", key=f"up_{idx}_{cause_idx}"):
                                st.session_state.feedback_states[feedback_key] = "positive"
                                feedback_data = {"timestamp": datetime.now().isoformat(), "feedback_type": "positive", "cause_rated": cause}
                                feedback_dir = project_root / "data" / "feedback"
                                feedback_dir.mkdir(exist_ok=True)
                                feedback_file = feedback_dir / "feedback_log.jsonl"
                                with open(feedback_file, "a") as f: f.write(json.dumps(feedback_data) + "\n")
                                str.write(f"Positive feedback logged for cause #{cause_idx+1} of drift #{idx+1}")
                                st.rerun()
                        with col2:
                            if st.button("👎", key=f"down_{idx}_{cause_idx}"):
                                st.session_state.feedback_states[feedback_key] = "negative"
                                feedback_data = {"timestamp": datetime.now().isoformat(), "feedback_type": "negative", "cause_rated": cause}
                                feedback_dir = project_root / "data" / "feedback"
                                feedback_dir.mkdir(exist_ok=True)
                                feedback_file = feedback_dir / "feedback_log.jsonl"
                                with open(feedback_file, "a") as f: f.write(json.dumps(feedback_data) + "\n")
                                str.write(f"Negative feedback logged for cause #{cause_idx+1} of drift #{idx+1}")
                                st.rerun()
    else:
        st.info("Select an event log and click 'Run Full Analysis' to begin.")

# Doc Viewer Tab
elif st.session_state.active_tab == "Document Viewer":
    if st.session_state.doc_to_view:
        display_document(st.session_state.doc_to_view)
    else:
        st.info("Click a document link to view it here.")

# --- Persistent Chatbot ---
if st.session_state.analysis_run_complete:
    if st.button("💬 Ask Follow-up Questions"):
        st.session_state.show_chat = not st.session_state.show_chat

    if st.session_state.show_chat:
        with st.container():
            st.write("---")
            st.subheader("Conversational Analysis")
            
            if not st.session_state.chat_history:
                st.session_state.chat_history.append(("assistant", "Hello! How can I help you with this analysis?"))

            for author, message in st.session_state.chat_history:
                avatar = "🧑‍💻" if author == "user" else "🤖"
                with st.chat_message(author, avatar=avatar):
                    st.markdown(message)

            if user_question := st.chat_input("Ask your question..."):
                st.session_state.chat_history.append(("user", user_question))
                
                with st.chat_message("user", avatar="🧑‍💻"):
                    st.markdown(user_question)
                    
                with st.chat_message("assistant", avatar="🤖"):
                    with st.spinner("Thinking..."):
                        current_state = {
                            "full_state_log": st.session_state.full_state_log,
                            "chat_history": st.session_state.chat_history,
                            "user_question": user_question
                        }
                        response = run_chatbot_agent(current_state)
                        st.session_state.chat_history = response.get('chat_history', [])
                        ai_answer = st.session_state.chat_history[-1][1]
                    
                    st.markdown(ai_answer)
                st.rerun()