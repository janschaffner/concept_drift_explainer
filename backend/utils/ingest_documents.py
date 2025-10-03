"""
This module provides a complete, multimodal data ingestion pipeline for populating
the Pinecone vector database.

It is designed to be run as a standalone script or to be called by the frontend
for on-the-fly document uploads. It can process a variety of file formats,
including text-based documents (.pdf, .docx, .txt), presentations (.pptx), and
standalone images (.png, .jpg). For multimodal files, it extracts both textual
content and uses a vision model to analyze and describe embedded images.

The script also handles the ingestion of a curated BPM glossary into a separate,
dedicated namespace within the same Pinecone index.
"""

import os
import sys
from pathlib import Path
from datetime import datetime
import logging
import io
import pandas as pd

# --- Path Correction ---
# This ensures that the script can correctly import modules from the 'backend' directory
# by adding the project's root directory to the system's path.
project_root = Path(__file__).resolve().parents[2]
sys.path.append(str(project_root))
# -----------------------

from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
from langchain.text_splitter import RecursiveCharacterTextSplitter
# Keep this even though it may be old
from langchain_community.document_loaders import UnstructuredFileLoader, PyPDFLoader
from langchain_openai import OpenAIEmbeddings
# Imports for multimodal processing
from backend.utils.image_analyzer import analyze_image_content
from pptx import Presentation

# --- Configuration ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Pinecone setup
PINECONE_DIMENSION = 1536
EMBEDDING_MODEL_NAME = "text-embedding-3-small" # This model has a 1536 dimension
# Define Namespace Constants
CONTEXT_NS = "context"
KB_NS      = "bpm-kb"


# Data path
DOCUMENTS_PATH = project_root / "frontend" / "static" / "documents"
# Create a temporary cache directory for extracted images
CACHE_DIR = project_root / "data" / "cache"
CACHE_DIR.mkdir(parents=True, exist_ok=True)


def get_timestamp_from_filename(filename: str) -> int:
    """Parses a YYYY-MM-DD date from a filename and returns a Unix timestamp.

    This function expects filenames to follow a convention like:
    '2023-10-27_document_name.pdf'. It extracts the date part and converts it
    to a Unix epoch timestamp, which is required for metadata filtering in Pinecone.

    Args:
        filename: The name of the file.

    Returns:
        An integer representing the Unix timestamp, or None if parsing fails.
    """
    try:
        date_str = filename.split('_')[0]
        dt_obj = datetime.strptime(date_str, '%Y-%m-%d')
        # Return timestamp as an integer (Unix epoch time)
        return int(dt_obj.timestamp())
    except (ValueError, IndexError):
        return None

# process_and_embed now accepts a namespace parameter
def process_and_embed(index, text_splitter, embedder, texts_to_embed: list, source_document_name: str, doc_timestamp: int, namespace: str):
    """Chunks, embeds, and upserts a list of texts into a Pinecone namespace.

    This is a generic helper function that takes raw text, splits it into
    manageable chunks, creates vector embeddings for each chunk, and then
    uploads them with their associated metadata to the specified namespace in
    the Pinecone index.

    Args:
        index: The initialized Pinecone index object.
        text_splitter: An initialized LangChain TextSplitter.
        embedder: An initialized LangChain embedding model.
        texts_to_embed: A list of raw text strings to be processed.
        source_document_name: The filename of the source document for metadata.
        doc_timestamp: The timestamp of the source document for metadata.
        namespace: The target Pinecone namespace (e.g., "context" or "bpm-kb").

    Returns:
        The number of vectors successfully ingested.
    """
    if not texts_to_embed:
        return 0
    
    # Use create_documents to handle a list of texts directly
    documents = text_splitter.create_documents(texts_to_embed)
    
    if not documents:
        logging.warning(f"Text splitting resulted in no chunks for {source_document_name}.")
        return 0
        
    texts = [chunk.page_content for chunk in documents]
    vectors = embedder.embed_documents(texts)
    
    vectors_to_upsert = []
    for i, chunk in enumerate(documents):
        # Create a more unique ID to avoid collisions between different documents
        vector_id = f"{Path(source_document_name).stem}_{hash(chunk.page_content)}_{i}"
        metadata = {
            "text": chunk.page_content,
            "source": source_document_name,
            "timestamp": doc_timestamp
        }
        vectors_to_upsert.append((vector_id, vectors[i], metadata))
        
    # Use the provided namespace for the upsert operation
    index.upsert(vectors=vectors_to_upsert, namespace=namespace)
    logging.info(f"  > Ingested {len(vectors_to_upsert)} vectors into namespace '{namespace}' from {source_document_name}.")
    return len(vectors_to_upsert)

# Function to specifically process the glossary file
def process_glossary_file(index, embedder):
    """Loads and ingests the BPM glossary into the 'bpm-kb' namespace.

    This function specifically processes the `bpm_glossary.csv` file, combining
    the term and its definition into a single text block for richer embeddings.
    It then upserts these into the dedicated 'bpm-kb' namespace to create a
    domain-specific knowledge layer for the retrieval agents.
    """
    glossary_path = project_root / "data" / "knowledge_base" / "bpm_glossary.csv"
    if not glossary_path.exists():
        logging.warning(f"BPM glossary not found at {glossary_path}. Skipping glossary ingestion.")
        return

    logging.info(f"--- Processing BPM Glossary ---")
    try:
        # Load the glossary from the CSV file.
        df = pd.read_csv(glossary_path)
        # Combine term and definition into a single text field for richer embeddings.
        df['text_to_embed'] = df['term'] + ": " + df['definition']
        
        texts = df['text_to_embed'].tolist()
        vectors = embedder.embed_documents(texts)
        
        vectors_to_upsert = []
        for i, row in df.iterrows():
            vector_id = f"bpm_kb_{i}"
            metadata = {
                "text": row['text_to_embed'],
                "source": "BPM Glossary",
                "timestamp": 0 # Timestamp is not relevant for the glossary
            }
            vectors_to_upsert.append((vector_id, vectors[i], metadata))
            
        # Upsert into the dedicated "bpm-kb" namespace using the constant
        index.upsert(vectors=vectors_to_upsert, namespace=KB_NS)
        logging.info(f"Successfully ingested {len(vectors_to_upsert)} terms into the '{KB_NS}' namespace.")

    except Exception as e:
        logging.error(f"Error processing glossary file: {e}")

# --- Main function to process context documents ---
def process_context_files(files_to_process: list, index, embedder, text_splitter):
    """Processes a list of documents, handling various file types.

    This is the main processing function for context documents. It iterates
    through a list of file paths, determines the file type, and uses the
    appropriate loader or utility to extract all textual and visual content.
    The extracted content is then passed to the `process_and_embed` helper to
    be ingested into the 'context' namespace.
    """
    logging.info(f"--- Processing {len(files_to_process)} Context Documents ---")
    
    total_vectors_ingested = 0
    for doc_path in files_to_process:
        logging.info(f"Processing document: {doc_path.name}")
        
        doc_timestamp = get_timestamp_from_filename(doc_path.name)
        if not doc_timestamp:
            logging.warning(f"Could not parse timestamp from '{doc_path.name}'. Skipping file.")
            continue
        
        # This list will aggregate all text extracted from a single file.
        texts_to_embed = []
        file_suffix = doc_path.suffix.lower()

        # Branch 1: Handle standalone image files.
        if file_suffix in [".png", ".jpg", ".jpeg"]:
            description = analyze_image_content(doc_path)
            if "Error" not in description:
                texts_to_embed.append(description)

        # Branch 2: Handle PowerPoint files, extracting both text and images.
        elif file_suffix == ".pptx":
            try:
                prs = Presentation(doc_path)
                for i, slide in enumerate(prs.slides):
                    # First, extract all text from shapes on the slide.
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts_to_embed.append(shape.text)
                    # Second, extract and analyze all images on the slide.
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            image = shape.image
                            # Save image to a temporary file to be analyzed.
                            temp_image_path = CACHE_DIR / f"temp_{image.sha1}.{image.ext}"
                            with open(temp_image_path, "wb") as f:
                                f.write(image.blob)
                            
                            image_description = analyze_image_content(temp_image_path)
                            if "Error" not in image_description:
                                # Add context to the description for better understanding.
                                full_description = f"Description of an image from slide {i+1} of '{doc_path.name}': {image_description}"
                                texts_to_embed.append(full_description)
                            os.remove(temp_image_path) # Clean up temp file.
            except Exception as e:
                logging.error(f"Error processing PowerPoint file {doc_path.name}: {e}")
                continue
        
        # Branch 3: Handle PDF files with the robust PyPDFLoader.
        elif file_suffix == ".pdf":
            try:
                loader = PyPDFLoader(str(doc_path))
                documents = loader.load_and_split()
                texts_to_embed.extend([doc.page_content for doc in documents])
            except Exception as e:
                logging.error(f"Error processing PDF file with PyPDFLoader {doc_path.name}: {e}")
                continue
        
        # Branch 4: Use a general loader for other text-based files (.docx, .txt).
        else:
            try:
                loader = UnstructuredFileLoader(str(doc_path))
                documents = loader.load()
                texts_to_embed.extend([doc.page_content for doc in documents])
            except Exception as e:
                logging.error(f"Error processing document with UnstructuredFileLoader {doc_path.name}: {e}")
                continue

        # After extracting all text, embed and upsert it into the 'context' namespace.            
        if texts_to_embed:
            # Call process_and_embed with the "context" namespace using the constant
            count = process_and_embed(index, text_splitter, embedder, texts_to_embed, doc_path.name, doc_timestamp, namespace=CONTEXT_NS)
            total_vectors_ingested += count
        else:
            logging.warning(f"No text or valid images could be extracted from '{doc_path.name}'.")
    
    return total_vectors_ingested

def initialize_ingestion_backend():
    """Initializes and returns all necessary backend resources for ingestion.

    This function centralizes the setup logic for Pinecone, the OpenAI embedder,
    and the LangChain text splitter. It checks if the Pinecone index exists and
    creates it if necessary.

    Returns:
        A tuple containing the initialized index, embedder, and text_splitter.
    """
    load_dotenv()
    api_key = os.getenv("PINECONE_API_KEY")
    pinecone_index_name = os.getenv("PINECONE_INDEX_NAME")

    if not all([api_key, pinecone_index_name]):
        raise ValueError("PINECONE_API_KEY or PINECONE_INDEX_NAME not found in .env file.")

    pc = Pinecone(api_key=api_key)

    if pinecone_index_name not in pc.list_indexes().names():
        logging.info(f"Index '{pinecone_index_name}' not found. Creating a new one...")
        pc.create_index(
            name=pinecone_index_name,
            dimension=PINECONE_DIMENSION,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )
        logging.info(f"Index '{pinecone_index_name}' created successfully.")
    else:
        logging.info(f"Connected to existing index: '{pinecone_index_name}'")

    index = pc.Index(pinecone_index_name)
    embedder = OpenAIEmbeddings(model=EMBEDDING_MODEL_NAME)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    logging.info("Backend ingestion resources initialized.")
    return index, embedder, text_splitter



if __name__ == "__main__":
    """
    Main execution block to run the script in standalone mode.

    This will initialize the backend, discover all supported document types in
    the `DOCUMENTS_PATH`, and process them. It will then process the BPM
    glossary file separately.
    """
    logging.info("--- Running Ingestion Script in Standalone Mode ---")
    
    # 1. Initialize all backend resources using the new reusable function
    index, embedder, text_splitter = initialize_ingestion_backend()
    
    # 2. Discover all supported files in the documents directory.
    all_files_to_process = (
        list(DOCUMENTS_PATH.glob("*.pdf")) +
        list(DOCUMENTS_PATH.glob("*.pptx")) +
        list(DOCUMENTS_PATH.glob("*.png")) +
        list(DOCUMENTS_PATH.glob("*.jpg")) +
        list(DOCUMENTS_PATH.glob("*.docx")) +
        list(DOCUMENTS_PATH.glob("*.txt"))
    )

    if all_files_to_process:
        # 3. Process the discovered context documents
        process_context_files(all_files_to_process, index, embedder, text_splitter)
    else:
        logging.warning(f"No context documents found in {DOCUMENTS_PATH} to process.")

    # 4. Process the glossary file
    process_glossary_file(index, embedder)

    logging.info("\n--- Full Ingestion Process Complete ---")